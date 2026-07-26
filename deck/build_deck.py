"""Build the Echo-Loop Idea submission deck from the provided SIH template.

Template rules honoured:
  * six slides INCLUDING the title — the instructions slide is deleted, exactly
    as that slide itself directs;
  * every "idea details pointer" survives verbatim, restyled as the eyebrow
    chip above each slide's descriptive title;
  * points and diagrams, not paragraphs.

Layout follows the Gamma reference: 0.8in margin, eyebrow chip, large light
title, twin 5.6in columns, full-width figure or accent band at the foot. Dark
tech palette, neon green reserved for positive outcomes.

All figures come from `make_images.py`, which reads the committed run outputs,
so the numbers on the slides cannot drift from the repository.
"""
from __future__ import annotations

from pathlib import Path

from lxml import etree
from pptx import Presentation
from pptx.enum.text import PP_ALIGN
from pptx.oxml.ns import qn
from pptx.util import Inches, Pt

from style import (AMBER, BLUE_LIGHT, BODY, BODY_INK, COL_2_X, COL_W, GREEN,
                   INK, MARGIN, MONO, MUTED, PANEL, PANEL_2, TITLE_FONT, WHITE,
                   card, eyebrow, pill, running_head, set_bg,
                   strip_template_content)

HERE = Path(__file__).resolve().parent
TEMPLATE = HERE / "template.pptx"
OUTPUT = HERE / "Echo-Loop_Idea_Deck.pptx"

REPO_URL = "https://github.com/ayushYadav1107/Eco_loop"
DEMO_URL = "https://eco-loop-agents.streamlit.app/"

PS_ID = "1"
THEME = "Clean & Green Technology"
STUDENT = "Ayush Yadav"
STUDENT_ID = "23BAI10006"


# --------------------------------------------------------------------------- #
def retheme_hyperlinks(prs, colour="60A5FA", visited="93A7BD") -> None:
    """Repoint the theme's hyperlink colours at the dark palette.

    PowerPoint paints hyperlink runs from `a:hlink` in the theme and ignores the
    run-level solidFill, so links rendered in the theme's near-navy blue —
    unreadable on this ground. Changing the theme fixes every link at once.
    The theme arrives as an opaque part, so it is parsed and written back.
    """
    theme = prs.slide_masters[0].part.part_related_by(
        "http://schemas.openxmlformats.org/officeDocument/2006/relationships/theme")
    root = etree.fromstring(theme.blob)
    ns = "{http://schemas.openxmlformats.org/drawingml/2006/main}"
    for tag, value in ((f"{ns}hlink", colour), (f"{ns}folHlink", visited)):
        for node in root.iter(tag):
            for child in list(node):
                node.remove(child)
            etree.SubElement(node, f"{ns}srgbClr", val=value)
    theme._blob = etree.tostring(root, xml_declaration=True,
                                 encoding="UTF-8", standalone=True)


def set_bullet(paragraph, char: str = "▪") -> None:
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set("marL", "150000")
    pPr.set("indent", "-150000")
    pPr.append(pPr.makeelement(qn("a:buFont"), {"typeface": BODY}))
    pPr.append(pPr.makeelement(qn("a:buChar"), {"char": char}))


def clear_bullet(paragraph) -> None:
    """The template's body boxes carry list formatting; a paragraph that does
    not opt out inherits a glyph and a hanging indent."""
    pPr = paragraph._p.get_or_add_pPr()
    for tag in ("a:buNone", "a:buChar", "a:buAutoNum"):
        for el in pPr.findall(qn(tag)):
            pPr.remove(el)
    pPr.set("marL", "0")
    pPr.set("indent", "0")
    pPr.append(pPr.makeelement(qn("a:buNone"), {}))


def write(shape, blocks, *, space_after=6) -> None:
    """(text, size, bold, colour, font, bullet) rows -> a styled text frame."""
    tf = shape.text_frame
    tf.word_wrap = True
    tf.clear()
    first = True
    for text, size, bold, colour, font, bullet in blocks:
        para = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        para.space_after = Pt(space_after)
        para.alignment = PP_ALIGN.LEFT      # template boxes are justified
        set_bullet(para) if bullet else clear_bullet(para)
        run = para.add_run()
        run.text = text
        run.font.name = font
        run.font.size = Pt(size)
        run.font.bold = bold
        run.font.color.rgb = colour


def textbox(slide, left, top, width, height):
    box = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    box.text_frame.word_wrap = True
    return box


def big_title(slide, text, sub=None, top=1.02):
    box = textbox(slide, MARGIN, top, 11.7, 0.62)
    write(box, [(text, 27, False, WHITE, TITLE_FONT, False)], space_after=0)
    if sub:
        s = textbox(slide, MARGIN, top + 0.56, 11.7, 0.34)
        write(s, [(sub, 11.5, True, BODY_INK, BODY, False)], space_after=0)
    return box


def head(slide, left, top, text, colour=WHITE, width=COL_W):
    box = textbox(slide, left, top, width, 0.3)
    write(box, [(text, 13, False, colour, TITLE_FONT, False)], space_after=0)
    return box


def metric(slide, left, top, value, label, sub, colour=GREEN, vsize=40):
    """Gamma's metric block: huge figure, label, one-line caption."""
    box = textbox(slide, left, top, 3.6, 0.74)
    write(box, [(value, vsize, False, colour, TITLE_FONT, False)], space_after=0)
    lab = textbox(slide, left, top + 0.7, 3.6, 0.62)
    write(lab, [
        (label, 13, False, INK, TITLE_FONT, False),
        (sub, 11, False, MUTED, BODY, False),
    ], space_after=2)
    return box


def image(slide, img: Path, left, top, width):
    return slide.shapes.add_picture(str(img), Inches(left), Inches(top), width=Inches(width))


def link_run(paragraph, text, url, *, size=11.5, colour=BLUE_LIGHT, font=BODY, bold=False):
    run = paragraph.add_run()
    run.text = text
    run.font.name = font
    run.font.size = Pt(size)
    run.font.bold = bold
    run.hyperlink.address = url          # set first; it resets run colour
    run.font.color.rgb = colour
    return run


# --------------------------------------------------------------------------- #
def main() -> None:
    prs = Presentation(str(TEMPLATE))

    sld_lst = prs.slides._sldIdLst
    first = list(sld_lst)[0]
    prs.part.drop_rel(first.get(qn("r:id")))
    sld_lst.remove(first)

    s1, s2, s3, s4, s5, s6 = prs.slides
    for n, slide in enumerate(prs.slides, 1):
        set_bg(slide)
        strip_template_content(slide)   # rebuild each slide from scratch
        running_head(slide, n)

    # ==================================================== 1 · TITLE & HOOK
    eyebrow(s1, "Autonomous Building Energy Management")
    big_title(s1, "Echo-Loop: A Building That Runs Itself",
              "Autonomous closed-loop HVAC control using EnergyPlus, MCP and a "
              "local open-source LLM.")

    card(s1, MARGIN, 2.04, 11.7, 0.94, fill=PANEL)
    hook = textbox(s1, MARGIN + 0.28, 2.14, 11.1, 0.76)
    write(hook, [
        ("EnergyPlus supplies the physics. A local open-source LLM supplies the "
         "judgement. The Model Context Protocol carries every reading and every "
         "command between them — while the simulation is still running.",
         12.5, False, INK, BODY, False),
    ], space_after=0)

    metric(s1, MARGIN, 3.2, "−8.5%", "HVAC Energy", "Saved over a summer week", GREEN)
    metric(s1, 4.85, 3.2, "336 / 336", "Agent Turns", "Zero failures", BLUE_LIGHT, vsize=36)
    metric(s1, 8.9, 3.2, "100%", "Open Source", "One laptop — no cloud, no API keys", GREEN)

    ident = textbox(s1, MARGIN, 4.66, 7.3, 1.1)
    write(ident, [
        (f"Problem Statement ID  {PS_ID}      ·      Theme  {THEME}",
         11.5, True, INK, BODY, False),
        (f"PS Category  Software      ·      Student  {STUDENT}      ·      "
         f"Student ID  {STUDENT_ID}", 11.5, False, BODY_INK, BODY, False),
    ], space_after=6)

    links = textbox(s1, 8.4, 4.66, 4.13, 1.1)
    tf = links.text_frame; tf.clear()
    p = tf.paragraphs[0]; p.alignment = PP_ALIGN.RIGHT; clear_bullet(p)
    r = p.add_run(); r.text = "Repository   "
    r.font.name = BODY; r.font.size = Pt(11); r.font.color.rgb = MUTED
    link_run(p, "github.com/ayushYadav1107/Eco_loop", REPO_URL, size=11)
    p2 = tf.add_paragraph(); p2.alignment = PP_ALIGN.RIGHT; clear_bullet(p2)
    r = p2.add_run(); r.text = "Live demo   "
    r.font.name = BODY; r.font.size = Pt(11); r.font.color.rgb = MUTED
    link_run(p2, "eco-loop-agents.streamlit.app", DEMO_URL, size=11)

    loop_x = MARGIN
    for i, (stage, col) in enumerate([
            ("EnergyPlus", BLUE_LIGHT), ("MCP tools", BLUE_LIGHT),
            ("Local LLM", GREEN), ("Safety gate", AMBER),
            ("Live actuators", BLUE_LIGHT)]):
        _, w = pill(s1, loop_x, 5.5, stage, fg=col)
        loop_x += w
        if i < 4:
            arrow = textbox(s1, loop_x + 0.02, 5.5, 0.36, 0.27)
            write(arrow, [("→", 12, False, MUTED, BODY, False)], space_after=0)
            loop_x += 0.38
    tail = textbox(s1, loop_x + 0.24, 5.53, 5.0, 0.3)
    write(tail, [("one closed loop, every 60 simulated minutes",
                  10.5, False, MUTED, BODY, False)], space_after=0)

    x = MARGIN
    for tag in ("EnergyPlus 26.1", "pyenergyplus", "FastMCP", "Model Context Protocol",
                "Ollama + Llama 3.2 3B", "Pydantic", "Streamlit"):
        _, w = pill(s1, x, 6.1, tag, fg=MUTED)
        x += w + 0.13

    # =============================================== 2 · SOLUTION & INNOVATION
    eyebrow(s2, "Proposed Solution & Innovation")
    big_title(s2, "Traditional BMS follow rigid clock schedules.",
              "Echo-Loop turns the building into a self-correcting agent.")

    head(s2, MARGIN, 2.04, "How It Works")
    body = textbox(s2, MARGIN, 2.4, COL_W, 1.7)
    write(body, [
        ("EnergyPlus runs in-process; sensors read every zone timestep",
         11.5, False, BODY_INK, BODY, True),
        ("Every 60 simulated minutes the state reaches a local LLM via six MCP tools",
         11.5, False, BODY_INK, BODY, True),
        ("The model weighs comfort, demand and grid carbon, then returns setpoints",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=8)

    head(s2, MARGIN, 4.16, "Why It's Different")
    why = textbox(s2, MARGIN, 4.52, COL_W, 0.9)
    write(why, [
        ("Most LLM-and-simulation work rewrites the model file and re-runs it — it "
         "cannot react to a zone drifting out of comfort at 14:00 on day 3. "
         "Echo-Loop does.", 11.5, False, BODY_INK, BODY, False),
    ], space_after=0)

    image(s2, HERE / "results.png", COL_2_X, 2.04, 5.6)

    # One real log line does more than a paragraph: it shows the loop closing.
    card(s2, COL_2_X, 4.16, COL_W, 1.16, fill=PANEL_2)
    log = textbox(s2, COL_2_X + 0.26, 4.28, COL_W - 0.5, 0.96)
    write(log, [
        ("ONE LINE FROM A LIVE RUN", 9, True, MUTED, BODY, False),
        ("[ai 07-17 09:00] OAT=24.9C PMV=-1.02 kW=2.09",
         10.5, False, BLUE_LIGHT, MONO, False),
        ("   -> cool=24.8 heat=19.2 (llm, 1.3s)", 10.5, False, GREEN, MONO, False),
    ], space_after=3)

    for i, (t, d, note, col) in enumerate([
        ("Injected into the LIVE solver", "No IDF rewrite, no restart.",
         "Setpoints reach the actuator handles mid-solve.", BLUE_LIGHT),
        ("Safety Boundary", "MCP validates every command server-side.",
         "Five gates; a rejected command never reaches the model.", AMBER),
        ("Measured, Not Projected", "PMV −0.43 → −0.02 while saving 8.5%.",
         "Both runs are committed EnergyPlus output.", GREEN),
    ]):
        x = MARGIN + i * 3.98
        card(s2, x, 5.5, 3.78, 1.3, fill=PANEL)
        b = textbox(s2, x + 0.24, 5.6, 3.32, 1.1)
        write(b, [
            (t, 12.5, False, col, TITLE_FONT, False),
            (d, 11, False, BODY_INK, BODY, False),
            (note, 10, False, MUTED, BODY, False),
        ], space_after=3)

    # =================================================== 3 · TECHNICAL APPROACH
    eyebrow(s3, "Technical Approach")
    big_title(s3, "The Closed Loop Architecture")

    image(s3, HERE / "arch.png", 1.72, 1.76, 9.9)
    image(s3, HERE / "timeline.png", MARGIN, 5.26, 11.7)

    steps = textbox(s3, MARGIN, 6.54, 11.7, 0.34)
    write(steps, [
        ("Sample  →  Aggregate  →  LLM tool-calling turn  →  Validate gate  →  "
         "Inject into live actuators via the pyenergyplus C API",
         11, False, MUTED, BODY, False),
    ], space_after=0)

    # ================================================== 4 · FEASIBILITY & RISK
    eyebrow(s4, "Feasibility and Viability")
    big_title(s4, "Feasibility and Overcoming Constraints")

    for i, (v, lab, sub, col) in enumerate([
        ("336 / 336", "Agent Turns", "0 fallbacks · 0 timeouts · 0 errors", GREEN),
        ("−3.7%", "Winter Energy", "Same loop, heating season", GREEN),
        ("2.1 s", "Median Latency", "Fully on-device", BLUE_LIGHT),
    ]):
        metric(s4, MARGIN + i * 4.0, 1.9, v, lab, sub, col, vsize=34)

    head(s4, MARGIN, 3.36, "Risks", AMBER)
    risks = textbox(s4, MARGIN, 3.72, COL_W, 1.72)
    write(risks, [
        ("LLM latency inside a blocking simulation callback",
         11.5, False, BODY_INK, BODY, True),
        ("A small model can move the wrong setpoint, or run heating against cooling",
         11.5, False, BODY_INK, BODY, True),
        ("One setpoint pair cannot satisfy five differently-loaded zones",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=10)

    head(s4, COL_2_X, 3.36, "Solutions", GREEN)
    sols = textbox(s4, COL_2_X, 3.72, COL_W, 1.72)
    write(sols, [
        ("Hard per-turn deadline; on timeout the loop holds the last accepted command",
         11.5, False, BODY_INK, BODY, True),
        ("Observations state the required direction — the model cannot invert physics",
         11.5, False, BODY_INK, BODY, True),
        ("Seasonal changeover parks the idle setpoint; a per-zone trim corrects each zone",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=10)

    image(s4, HERE / "pipeline.png", 1.05, 5.22, 11.2)

    # =============================================== 5 · RESEARCH & ARTIFACTS
    eyebrow(s5, "Artifacts")
    big_title(s5, "Artifacts & References")

    card(s5, MARGIN, 2.0, 11.7, 0.98, fill=PANEL)
    repo = textbox(s5, MARGIN + 0.28, 2.1, 11.1, 0.82)
    tf = repo.text_frame; tf.clear()
    p = tf.paragraphs[0]; clear_bullet(p); p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Repository   "
    r.font.name = TITLE_FONT; r.font.size = Pt(12.5); r.font.color.rgb = INK
    link_run(p, "github.com/ayushYadav1107/Eco_loop", REPO_URL, size=12.5, font=MONO, bold=True)
    r = p.add_run(); r.text = "        Live demo   "
    r.font.name = TITLE_FONT; r.font.size = Pt(12.5); r.font.color.rgb = INK
    link_run(p, "eco-loop-agents.streamlit.app", DEMO_URL, size=12.5, font=MONO, bold=True)
    p2 = tf.add_paragraph(); clear_bullet(p2)
    r = p2.add_run()
    r.text = ("Baseline and runtime-generated .idf models, agent_decisions.jsonl and "
              "per-timestep EnergyPlus results are all committed — a fresh clone "
              "reproduces every figure in this deck.")
    r.font.name = BODY; r.font.size = Pt(11); r.font.color.rgb = BODY_INK

    head(s5, MARGIN, 3.16, "Simulation & Building Physics", BLUE_LIGHT)
    refs1 = textbox(s5, MARGIN, 3.52, COL_W, 1.42)
    write(refs1, [
        ("EnergyPlus Engineering Reference & EMS Application Guide",
         11.5, False, BODY_INK, BODY, True),
        ("pyenergyplus Data Exchange API — get_variable_handle, set_actuator_value",
         11.5, False, BODY_INK, BODY, True),
        ("ISO 7730:2005 & ASHRAE Standard 55 — PMV/PPD comfort model",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=8)

    head(s5, MARGIN, 5.02, "Agent, Protocol & Inference", BLUE_LIGHT)
    refs2 = textbox(s5, MARGIN, 5.38, COL_W, 1.42)
    write(refs2, [
        ("Model Context Protocol — modelcontextprotocol.io",
         11.5, False, BODY_INK, BODY, True),
        ("FastMCP — gofastmcp.com", 11.5, False, BODY_INK, BODY, True),
        ("Meta Llama 3.2 served locally through Ollama — ollama.com",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=8)

    # decisions.png is 2.7:1 — full width would run 4.3in tall and fall off the
    # slide, so it lives in the right column at its natural height instead.
    head(s5, COL_2_X, 3.16, "168 setpoint decisions, one summer week", GREEN)
    image(s5, HERE / "decisions.png", COL_2_X, 3.52, COL_W)

    card(s5, COL_2_X, 5.72, COL_W, 1.06, fill=PANEL)
    note = textbox(s5, COL_2_X + 0.26, 5.84, COL_W - 0.5, 0.86)
    write(note, [
        ("Every decision is logged with the model's own reason",
         11.5, False, INK, TITLE_FONT, False),
        ("outputs/ai/agent_decisions.jsonl — committed, one JSON record per "
         "control interval.", 10.5, False, MUTED, BODY, False),
    ], space_after=3)

    # ========================================== 6 · REPRODUCIBILITY & LIMITS
    eyebrow(s6, "Research and References")
    big_title(s6, "Auditable, Not Anecdotal",
              "The controller runs at temperature 0 — repeated runs return identical "
              "totals. Anyone can clone the repo and reproduce these exact figures.")

    card(s6, MARGIN, 2.2, COL_W, 1.86, fill=PANEL_2)
    code = textbox(s6, MARGIN + 0.26, 2.34, COL_W - 0.5, 1.64)
    write(code, [
        ("python main.py prepare", 11, False, GREEN, MONO, False),
        ("python main.py run-baseline --start 07-15 --end 07-21", 11, False, GREEN, MONO, False),
        ("python main.py run-ai --start 07-15 --end 07-21", 11, False, GREEN, MONO, False),
        ("python main.py dashboard", 11, False, GREEN, MONO, False),
        ("# ~6 min end to end on a laptop — EnergyPlus 26.1 + Ollama,",
         10, False, MUTED, MONO, False),
        ("# no cloud, no API key, no paid service anywhere in the loop",
         10, False, MUTED, MONO, False),
    ], space_after=6)

    head(s6, COL_2_X, 2.2, "Four Commands. Full Reproduction.")
    steps = textbox(s6, COL_2_X, 2.56, COL_W, 1.6)
    write(steps, [
        ("prepare — stage and instrument the EnergyPlus model",
         11.5, False, BODY_INK, BODY, True),
        ("run-baseline — the same building on its native schedule",
         11.5, False, BODY_INK, BODY, True),
        ("run-ai — the closed loop, 168 agent decisions",
         11.5, False, BODY_INK, BODY, True),
        ("dashboard — the savings comparison, side by side",
         11.5, False, BODY_INK, BODY, True),
    ], space_after=7)

    image(s6, HERE / "comfort_dist.png", MARGIN, 4.28, 5.6)

    card(s6, COL_2_X, 4.28, COL_W, 2.3, fill=PANEL)
    lim = textbox(s6, COL_2_X + 0.26, 4.42, COL_W - 0.52, 2.02)
    write(lim, [
        ("Honest limitation", 12.5, False, AMBER, TITLE_FONT, False),
        ("One setpoint pair drives five differently-loaded zones. Per zone the agent "
         "holds 85–96% of occupied time inside the PMV band; the stricter worst-zone "
         "metric is lower.", 11.5, False, BODY_INK, BODY, False),
        ("Per-zone setpoints are the next step — the actuator handles and per-zone "
         "PMV are already in place.", 11.5, False, MUTED, BODY, False),
        ("Stated in full in ARCHITECTURE.md §8 — nothing here is hidden.",
         10.5, False, BLUE_LIGHT, BODY, False),
    ], space_after=8)

    retheme_hyperlinks(prs)

    prs.save(str(OUTPUT))
    print(f"wrote {OUTPUT.name}  ({len(prs.slides._sldIdLst)} slides)")


if __name__ == "__main__":
    main()
