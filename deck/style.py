"""Visual language for the Echo-Loop Idea deck.

Layout grammar is taken from the Gamma reference deck: a 0.8in margin, an
eyebrow chip naming the section, a large light title beneath it, twin 5.6in
columns, and a full-width accent band at the foot. The type scale is Gamma's
too — 27pt title, 13pt section head, 11.5pt body, 9pt eyebrow, 40pt metric.

Palette is light: a white ground with pale panels, a sustainability green
reserved for positive outcomes, a deep tech blue for structure, and a muted
orange for constraints and the heating setpoint. Every accent is darkened far
enough to clear 4.5:1 against both the ground and the panels — the dark theme's
neon green and vibrant blue sit near 1.6:1 and 3.7:1 on white and cannot be
reused as-is.

Font note. The reference uses Barlow, which is not installed here — PowerPoint
substitutes silently at export, so the PDF would not carry it either. Segoe UI
(the closest installed grotesque) is used instead, named in one place below.
"""
from __future__ import annotations

from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# --------------------------------------------------------------------------- #
# palette — light. Contrast against the white ground is given per accent.
# --------------------------------------------------------------------------- #
GROUND = RGBColor(0xFF, 0xFF, 0xFF)     # slide background
PANEL = RGBColor(0xF1, 0xF5, 0xF9)      # raised card
PANEL_2 = RGBColor(0xE8, 0xEF, 0xF7)    # second-level card
HAIRLINE = RGBColor(0xCB, 0xD5, 0xE1)

TITLE_INK = RGBColor(0x0B, 0x12, 0x20)  # slide titles          18.9:1
INK = RGBColor(0x0F, 0x1B, 0x2D)        # primary text          16.6:1
BODY_INK = RGBColor(0x33, 0x45, 0x5C)   # body copy              9.4:1
MUTED = RGBColor(0x5A, 0x6B, 0x80)      # secondary              5.4:1

GREEN = RGBColor(0x0A, 0x7A, 0x52)      # positive outcomes ONLY 5.4:1
BLUE = RGBColor(0x1D, 0x4E, 0xD8)       # structure              6.7:1
BLUE_LIGHT = RGBColor(0x25, 0x63, 0xEB) # links, section heads   5.2:1
AMBER = RGBColor(0xC2, 0x41, 0x0C)      # constraints, heating   5.2:1

# Requested Barlow -> installed substitute. Swap here after installing Barlow.
TITLE_FONT = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

MARGIN = 0.8
COL_W = 5.6
COL_2_X = 6.93


# --------------------------------------------------------------------------- #
def set_bg(slide, colour: RGBColor = GROUND) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = colour


def send_to_back(shape) -> None:
    sp = shape._element
    tree = sp.getparent()
    tree.remove(sp)
    tree.insert(2, sp)


def card(slide, left, top, width, height, *, fill=PANEL, line=None, radius=0.04):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height))
    shape.adjustments[0] = radius
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill
    if line is None:
        shape.line.fill.background()
    else:
        shape.line.color.rgb = line
        shape.line.width = Pt(1)
    shape.shadow.inherit = False
    shape.text_frame.text = ""
    send_to_back(shape)
    return shape


def pill(slide, left, top, text, *, fg=GREEN, bg=None, fsize=9):
    """Gamma's tag pill: a rounded chip with small uppercase text."""
    width = 0.082 * len(text) + 0.30
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE, Inches(left), Inches(top),
        Inches(width), Inches(0.27))
    shape.adjustments[0] = 0.5
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg or PANEL_2
    shape.line.fill.background()
    shape.shadow.inherit = False
    tf = shape.text_frame
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.CENTER
    run = para.add_run()
    run.text = text.upper()
    run.font.name = BODY
    run.font.size = Pt(fsize)
    run.font.bold = True
    run.font.color.rgb = fg
    return shape, width


def strip_template_content(slide) -> None:
    """Remove every shape the SIH template ships on the slide.

    The template's body boxes carry 28-32pt prompt text ("Problem Statement
    ID -", "Student Name ...") and their own list/justification formatting. The
    deck now composes each slide from scratch, so leaving them behind renders
    that prompt text as giant ghost copy underneath the real content — which is
    exactly what happened on the first build. Everything is rebuilt explicitly.
    """
    for sh in list(slide.shapes):
        sh._element.getparent().remove(sh._element)


def eyebrow(slide, text, *, colour=GREEN):
    """Gamma's section chip: small uppercase label above the title.

    The SIH "idea details pointer" for each slide is carried here verbatim, so
    the required section names survive even though the layout is rebuilt.
    """
    box = slide.shapes.add_textbox(Inches(MARGIN), Inches(0.56), Inches(7.0), Inches(0.3))
    tf = box.text_frame
    tf.word_wrap = False
    tf.margin_left = tf.margin_top = tf.margin_bottom = 0
    para = tf.paragraphs[0]
    para.alignment = PP_ALIGN.LEFT
    run = para.add_run()
    run.text = text.upper()
    run.font.name = BODY
    run.font.size = Pt(10)
    run.font.bold = True
    run.font.color.rgb = colour
    return box


def strip_bottom_chrome(slide) -> None:
    """Delete the template's footer and slide-number placeholders outright.

    Emptying them leaves invisible shapes that still occupy the layout and
    still surface in any later audit; the brief asks for them gone.
    """
    for sh in list(slide.shapes):
        if "Footer" in sh.name or "Slide Number" in sh.name or "Date" in sh.name:
            sh._element.getparent().remove(sh._element)


def running_head(slide, number: int) -> None:
    """'Echo-Loop' on the left, slide number on the right. No box, no rule."""
    left = slide.shapes.add_textbox(Inches(MARGIN), Inches(6.94), Inches(4.0), Inches(0.32))
    p = left.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.LEFT
    r = p.add_run(); r.text = "Echo-Loop"
    r.font.name = BODY; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = MUTED

    right = slide.shapes.add_textbox(Inches(8.5), Inches(6.94), Inches(4.03), Inches(0.32))
    p = right.text_frame.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run(); r.text = str(number)
    r.font.name = BODY; r.font.size = Pt(10); r.font.bold = True
    r.font.color.rgb = GREEN
